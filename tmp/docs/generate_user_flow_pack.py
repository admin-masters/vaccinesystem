#!/usr/bin/env python3
from __future__ import annotations

import json
import re
import zipfile
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
MANUAL_DIR = ROOT / "docs" / "product-user-flows"
ASSET_DIR = MANUAL_DIR / "assets"
DECK_DIR = ROOT / "output" / "doc" / "user-flow-decks"
QA_REPORT = DECK_DIR / "qa-report.txt"


COLORS = {
    "ink": RGBColor(31, 41, 55),
    "muted": RGBColor(71, 85, 105),
    "teal": RGBColor(15, 118, 110),
    "teal_light": RGBColor(204, 251, 241),
    "coral": RGBColor(249, 115, 22),
    "coral_light": RGBColor(255, 237, 213),
    "surface": RGBColor(248, 250, 252),
    "line": RGBColor(226, 232, 240),
}

STATUS_TEXT = "Validated against the local demo build on April 9, 2026. Documentation reflects live behavior; mismatches with older expectations are noted."


def asset_path(slug: str, filename: str) -> str:
    return f"docs/product-user-flows/assets/{slug}/{filename}"


WORKFLOWS: list[dict[str, Any]] = [
    {
        "number": 1,
        "group": "Platform Overview",
        "slug": "platform-overview-and-role-map",
        "title": "Platform Overview and Role Map",
        "purpose": "Orient trainers, internal teams, and new stakeholders to the product’s end-to-end operating model before they dive into role-specific workflows.",
        "primary_user": "Internal trainers, product owners, partner leads, and new project team members",
        "entry_point": "Public home page `/` plus direct workflow URLs for parents, doctors, and partners",
        "summary": "The product is a pediatric vaccination system connecting partner onboarding, doctor registration, patient record management, parent self-service, reminders, and education content. The live app currently starts from a doctor/admin-oriented home page, while parent flows start from direct routes such as `/add/` and `/update/`.",
        "diagram": """```mermaid
flowchart LR
    A["Partner / Admin"] --> B["Partner link or doctor registration support"]
    B --> C["Doctor / Clinic Portal"]
    C --> D["Add Patient or Update Vaccination Record"]
    D --> E["Parent Card Access"]
    C --> F["Reminders Dashboard"]
    F --> E
    E --> G["Vaccination History and Education"]
```""",
        "success_criteria": [
            "The trainer can explain which role enters at which route.",
            "The audience understands how partner setup, doctor usage, parent access, reminders, and education fit together.",
            "Known live-product mismatches are called out before role training begins.",
        ],
        "related": [
            "02-admin-partner-provisioning-and-field-rep-upload.md",
            "03-doctor-self-registration-and-portal-access.md",
            "08-parent-self-service-card-history-and-share-link.md",
        ],
        "notes": [
            "The home page currently exposes doctor registration and a Django admin link, not the parent actions described in the original brief.",
            "Parent workflows are live, but trainers should take users directly to `/add/` or `/update/` until the home page is redesigned.",
        ],
        "steps": [
            {
                "title": "Open the live landing page",
                "user_does": "Open the public home page in a browser.",
                "user_sees": "A simple landing page with Doctor Registration and Field Partner (Admin) actions.",
                "why": "This establishes the real starting point the product exposes today.",
                "result": "The trainer can explain that the live landing page is provider/admin oriented.",
                "notes": "Call out that this differs from the original expectation of a parent-first landing page.",
                "screenshot_path": asset_path("platform-overview-and-role-map", "01-home-overview.png"),
                "caption": "Live home page showing the current public entry points",
                "what_to_show": "The home page with the doctor registration and admin actions visible.",
                "require_asset": True,
            },
            {
                "title": "Map each role to its real entry point",
                "user_does": "Review the product with role-based starting URLs: admin/partner, doctor, and parent.",
                "user_sees": "A product where some roles start from the home page and others start from direct routes or shared links.",
                "why": "Role clarity prevents confusion during onboarding and training delivery.",
                "result": "Users understand that parents, doctors, and partner/admin users do not all start from the same page.",
                "notes": "Parent flows are best demonstrated from `/add/`, `/update/`, or a share link.",
                "screenshot_path": asset_path("platform-overview-and-role-map", "01-home-overview.png"),
                "caption": "Home page reused as the anchor for the role-entry explanation",
                "what_to_show": "The live landing page while the trainer overlays the role map verbally or in the deck.",
                "require_asset": True,
            },
            {
                "title": "Explain the end-to-end handoff chain",
                "user_does": "Walk through how partner setup enables doctor onboarding, which then enables patient registration, reminders, and parent access.",
                "user_sees": "A connected service flow rather than isolated screens.",
                "why": "This is the mental model that makes the rest of the workflow pack coherent.",
                "result": "The trainer can position each detailed workflow deck inside one shared operating model.",
                "notes": "Use the role map slide in the deck for this step; there is no single live UI screen for the whole chain.",
                "screenshot_path": asset_path("platform-overview-and-role-map", "02-role-map-diagram.png"),
                "caption": "Role map diagram for the end-to-end service flow",
                "what_to_show": "A simple visual showing partner/admin -> doctor/clinic -> parent/patient progression.",
                "require_asset": False,
            },
        ],
    },
    {
        "number": 2,
        "group": "Admin / Partner",
        "slug": "admin-partner-provisioning-and-field-rep-upload",
        "title": "Admin Partner Provisioning and Field Rep Upload",
        "purpose": "Show how a partner administrator opens the partner workspace, creates a partner profile, uploads field representatives, and generates a doctor registration link.",
        "primary_user": "Partner admin, field operations lead, or internal business operations owner",
        "entry_point": "Admin access gate at `/admin/access/`, followed by `/partners/new/`",
        "summary": "The live product provides a lightweight admin password gate for partner setup. After access is granted, the admin can create a partner, upload field rep CSV data, and obtain a doctor registration link to circulate.",
        "success_criteria": [
            "The admin can reach the partner workspace.",
            "A partner record is created with field reps loaded from CSV.",
            "A working doctor registration link is produced for onward sharing.",
        ],
        "related": [
            "04-partner-led-doctor-onboarding.md",
            "03-doctor-self-registration-and-portal-access.md",
        ],
        "notes": [
            "The public home page links to Django admin, but the actual partner provisioning flow is the custom `/admin/access/` -> `/partners/new/` path.",
        ],
        "steps": [
            {
                "title": "Pass the admin access gate",
                "user_does": "Open the admin access page and enter the configured admin password.",
                "user_sees": "A focused password gate that protects the partner provisioning workspace.",
                "why": "This is the first gate that separates general users from partner administrators.",
                "result": "The session is authorized for partner creation.",
                "notes": "In the local demo build, the admin quick password is derived from the local secret key.",
                "screenshot_path": asset_path("admin-partner-provisioning-and-field-rep-upload", "01-admin-access.png"),
                "caption": "Admin access gate before partner setup",
                "what_to_show": "The password prompt at `/admin/access/`.",
                "require_asset": True,
            },
            {
                "title": "Create the partner workspace",
                "user_does": "Open the partner creation screen after access is granted.",
                "user_sees": "A form with Partner Name and Field Reps CSV upload fields.",
                "why": "This is where the whitelabel / partner shell and field rep import are initiated.",
                "result": "The admin is ready to create the partner record.",
                "notes": "The field rep CSV expects `rep_code` and `full_name` columns.",
                "screenshot_path": asset_path("admin-partner-provisioning-and-field-rep-upload", "02-partner-create-form.png"),
                "caption": "Partner creation form with CSV upload",
                "what_to_show": "The Create Partner & Upload Field Reps form.",
                "require_asset": True,
            },
            {
                "title": "Generate and share the doctor registration link",
                "user_does": "Submit the partner form with the field rep CSV.",
                "user_sees": "A success message showing the generated doctor registration link.",
                "why": "This link is the handoff artifact used by field teams and partner operators to onboard doctors.",
                "result": "The partner exists and its registration link can be shared.",
                "notes": "Treat the link as the official onboarding URL for partner-led doctor registration.",
                "screenshot_path": asset_path("admin-partner-provisioning-and-field-rep-upload", "03-partner-created.png"),
                "caption": "Partner creation success state with the generated onboarding link",
                "what_to_show": "The success message after partner creation, including the registration link text.",
                "require_asset": True,
            },
        ],
    },
    {
        "number": 3,
        "group": "Doctor / Clinic",
        "slug": "doctor-self-registration-and-portal-access",
        "title": "Doctor Self Registration and Portal Access",
        "purpose": "Train a doctor or clinic lead to self-register, understand the expected WhatsApp handoff, and recognize the portal landing page they use after sign-in.",
        "primary_user": "Doctor or clinic lead registering directly",
        "entry_point": "Doctor self-registration form at `/doctor/register/`",
        "summary": "Doctors can self-register through a public registration form. After submission, the current product opens WhatsApp with a prefilled portal message, and the doctor then uses the portal link to reach the clinic workspace.",
        "success_criteria": [
            "The doctor can complete the registration form.",
            "The doctor understands that submission hands off to WhatsApp with the portal link.",
            "The doctor recognizes the portal home once authenticated.",
        ],
        "related": [
            "05-doctor-portal-add-patient-and-share-link.md",
            "06-doctor-portal-update-card-and-profile.md",
        ],
        "notes": [
            "Production behavior expects Google/Gmail sign-in for portal use.",
            "The local documentation build captures the portal home using a demo session because local OAuth credentials are placeholders.",
        ],
        "steps": [
            {
                "title": "Open the doctor registration form",
                "user_does": "Browse to the self-registration page.",
                "user_sees": "A detailed doctor registration form covering identity, clinic details, languages, contact numbers, and IMC number.",
                "why": "This is the doctor’s main onboarding form when no partner-issued link is involved.",
                "result": "The doctor can review the data required to register.",
                "notes": "The form enforces Gmail-format email addresses in the current build.",
                "screenshot_path": asset_path("doctor-self-registration-and-portal-access", "01-doctor-self-registration-form.png"),
                "caption": "Doctor self-registration form",
                "what_to_show": "The full doctor registration form with the mandatory fields visible.",
                "require_asset": True,
            },
            {
                "title": "Submit and hand off to WhatsApp",
                "user_does": "Complete the form and click Register.",
                "user_sees": "The browser is redirected to WhatsApp with a prefilled message containing the doctor portal link.",
                "why": "This is how the product currently distributes the personalized portal URL.",
                "result": "The doctor receives a reusable portal link message.",
                "notes": "The current implementation opens WhatsApp rather than using a server-side messaging gateway.",
                "screenshot_path": asset_path("doctor-self-registration-and-portal-access", "01-doctor-self-registration-form.png"),
                "caption": "Registration form immediately before submission",
                "what_to_show": "The completed form or the register action, since the next state is an external WhatsApp redirect.",
                "require_asset": True,
            },
            {
                "title": "Recognize the clinic portal home",
                "user_does": "Open the doctor portal link and sign in.",
                "user_sees": "A doctor portal home page with actions for Add-a-patient, Update a patient record, Reminders, and Edit doctor/clinic profile.",
                "why": "This is the control center the doctor returns to for daily clinic use.",
                "result": "The doctor can navigate to the main clinic workflows from one landing page.",
                "notes": "The documentation screenshots use a local demo session to bypass unavailable Google OAuth credentials.",
                "screenshot_path": asset_path("doctor-self-registration-and-portal-access", "02-doctor-self-portal-home.png"),
                "caption": "Doctor portal home after successful access",
                "what_to_show": "The doctor portal home with the primary action buttons visible.",
                "require_asset": True,
            },
        ],
    },
    {
        "number": 4,
        "group": "Admin / Partner",
        "slug": "partner-led-doctor-onboarding",
        "title": "Partner-Led Doctor Onboarding",
        "purpose": "Show how a partner-issued registration link adds field rep attribution to the doctor onboarding process and leads into the same portal experience.",
        "primary_user": "Field team member, partner coordinator, or doctor onboarding via partner link",
        "entry_point": "Partner registration URL at `/doctor/register/<token>/`",
        "summary": "Partner-led onboarding uses a tokenized doctor registration link. The form is similar to self-registration but adds field rep code and field rep name to attribute the onboarding to the partner team.",
        "success_criteria": [
            "The doctor can open the partner-specific registration form.",
            "The form captures the required field rep information.",
            "The doctor reaches the same portal destination after the WhatsApp handoff.",
        ],
        "related": [
            "02-admin-partner-provisioning-and-field-rep-upload.md",
            "03-doctor-self-registration-and-portal-access.md",
        ],
        "notes": [
            "The partner token is the mechanism that scopes the doctor to the partner.",
            "A valid field rep code is required for this path.",
        ],
        "steps": [
            {
                "title": "Open the partner-issued registration link",
                "user_does": "Open the doctor registration URL received from the partner/admin team.",
                "user_sees": "The doctor registration form with additional Field Rep Code and Field Rep Name inputs.",
                "why": "This route distinguishes partner-led onboarding from self-registration.",
                "result": "The doctor is in the correct registration context for partner attribution.",
                "notes": "If the token is invalid or expired, the product redirects back to the home page with an error.",
                "screenshot_path": asset_path("partner-led-doctor-onboarding", "01-partner-led-registration-form.png"),
                "caption": "Partner-led doctor registration form with field rep fields",
                "what_to_show": "The registration form including Field Rep Code and Field Rep Name.",
                "require_asset": True,
            },
            {
                "title": "Complete registration with field rep details",
                "user_does": "Enter doctor details plus the partner field rep code/name and submit the form.",
                "user_sees": "The same WhatsApp redirect behavior used for self-registration, now tied to the partner and field rep.",
                "why": "This preserves partner attribution while still delivering the portal link to the doctor.",
                "result": "The doctor is stored with the linked partner and field rep association.",
                "notes": "The live validation checks the field rep code against the selected partner.",
                "screenshot_path": asset_path("partner-led-doctor-onboarding", "01-partner-led-registration-form.png"),
                "caption": "Partner-led registration form before submission",
                "what_to_show": "The partner registration page prior to the WhatsApp redirect.",
                "require_asset": True,
            },
            {
                "title": "Confirm the partner doctor portal landing",
                "user_does": "Open the personalized portal link after sign-in.",
                "user_sees": "The same doctor portal home used in self-registration-based onboarding.",
                "why": "This confirms that partner-led onboarding still leads into the shared clinic operating workflow.",
                "result": "The partner-onboarded doctor can continue with routine clinic tasks.",
                "notes": "The downstream portal experience is shared; the main difference is the onboarding attribution.",
                "screenshot_path": asset_path("partner-led-doctor-onboarding", "02-partner-doctor-portal-home.png"),
                "caption": "Partner-onboarded doctor landing in the clinic portal",
                "what_to_show": "The partner doctor’s portal home after access.",
                "require_asset": True,
            },
        ],
    },
    {
        "number": 5,
        "group": "Doctor / Clinic",
        "slug": "doctor-portal-add-patient-and-share-link",
        "title": "Doctor Portal Add Patient and Share Link",
        "purpose": "Teach clinic staff how to register a patient, create the patient vaccination card, and understand the optional share-link handoff to the parent.",
        "primary_user": "Doctor or clinic staff member working inside the portal",
        "entry_point": "Doctor portal home `/d/<token>/` -> Add-a-patient",
        "summary": "From the portal, the doctor can add a patient by recording child details and the parent’s WhatsApp number. The portal supports a standard submit path and also shows a Register & Send to Patient action that opens WhatsApp Web with a shareable card link.",
        "success_criteria": [
            "The clinic can open the add-patient form.",
            "Submitting the form creates the child vaccination card.",
            "Staff understand that the optional send action opens WhatsApp Web with a share link.",
        ],
        "related": [
            "06-doctor-portal-update-card-and-profile.md",
            "08-parent-self-service-card-history-and-share-link.md",
        ],
        "notes": [
            "The product currently opens WhatsApp Web rather than sending through a business messaging gateway.",
        ],
        "steps": [
            {
                "title": "Open the add-patient form",
                "user_does": "Choose Add-a-patient from the portal home.",
                "user_sees": "The child registration form with Submit and Register & Send to Patient actions.",
                "why": "This is the clinic-side starting point for bringing a new child into the system.",
                "result": "The clinic user can record patient details and decide whether to also hand off a share link.",
                "notes": "The send button is only shown inside the doctor portal, not on the public parent add flow.",
                "screenshot_path": asset_path("doctor-portal-add-patient-and-share-link", "01-doctor-add-patient-form.png"),
                "caption": "Doctor portal add-patient form with both action buttons",
                "what_to_show": "The form and the Register & Send to Patient button.",
                "require_asset": True,
            },
            {
                "title": "Create the patient record",
                "user_does": "Submit the child details and parent WhatsApp number.",
                "user_sees": "The child’s vaccination card inside the doctor portal.",
                "why": "This confirms that the patient record and linked vaccination schedule were created successfully.",
                "result": "The doctor can immediately review and update the child card.",
                "notes": "The child is scoped to the doctor’s clinic for later update and reminder workflows.",
                "screenshot_path": asset_path("doctor-portal-add-patient-and-share-link", "02-doctor-card-after-add.png"),
                "caption": "Doctor view of the vaccination card after adding a patient",
                "what_to_show": "The doctor-facing vaccination card returned after patient creation.",
                "require_asset": True,
            },
            {
                "title": "Use the optional share-link path",
                "user_does": "Choose Register & Send to Patient instead of the standard submit action when appropriate.",
                "user_sees": "WhatsApp Web opens with a bilingual message and the parent share link.",
                "why": "This enables the clinic to hand off self-service access to the parent at the moment of registration.",
                "result": "The parent receives a direct verification link to the child card.",
                "notes": "The final messaging step happens outside the app in WhatsApp Web.",
                "screenshot_path": asset_path("doctor-portal-add-patient-and-share-link", "01-doctor-add-patient-form.png"),
                "caption": "The add-patient screen where the share-link handoff is initiated",
                "what_to_show": "The same add-patient form, highlighting the Register & Send to Patient button.",
                "require_asset": True,
            },
        ],
    },
    {
        "number": 6,
        "group": "Doctor / Clinic",
        "slug": "doctor-portal-update-card-and-profile",
        "title": "Doctor Portal Update Card and Profile",
        "purpose": "Train clinic users on finding an existing patient, opening the card, updating vaccine dates, and maintaining the doctor/clinic profile.",
        "primary_user": "Doctor or clinic staff member",
        "entry_point": "Doctor portal home -> Update a patient record, plus Edit doctor/clinic profile",
        "summary": "This workflow covers the day-to-day clinic maintenance path: look up a parent by WhatsApp number, select the correct child, open the vaccination card, and update clinic profile details as needed.",
        "success_criteria": [
            "The clinic can find patients by parent WhatsApp number.",
            "The correct child record is opened within the clinic scope.",
            "The doctor/clinic profile can be updated from the same portal.",
        ],
        "related": [
            "05-doctor-portal-add-patient-and-share-link.md",
            "07-reminders-and-education-workflows.md",
        ],
        "notes": [
            "Patient lookup is clinic-scoped in the doctor workflow, even if the parent has records elsewhere.",
        ],
        "steps": [
            {
                "title": "Search by parent WhatsApp number",
                "user_does": "Open Update a patient record and enter the parent’s WhatsApp number.",
                "user_sees": "A WhatsApp lookup form specific to the doctor portal.",
                "why": "This is the quickest way to retrieve a family’s children already linked to the clinic.",
                "result": "The system is ready to show the matching child list.",
                "notes": "The number lookup is based on the stored WhatsApp number.",
                "screenshot_path": asset_path("doctor-portal-update-card-and-profile", "01-doctor-update-lookup.png"),
                "caption": "Doctor-side lookup form for an existing patient",
                "what_to_show": "The clinic-scoped WhatsApp lookup screen.",
                "require_asset": True,
            },
            {
                "title": "Choose the correct child",
                "user_does": "Submit the lookup and review the matching children.",
                "user_sees": "A child-selection screen with actions for Open Vaccination Card, Full Schedule, and Send Reminders.",
                "why": "A single parent may have multiple children, so the clinic must choose the right record before editing anything.",
                "result": "The clinic can continue into the correct child workflow.",
                "notes": "Only children in the current clinic appear in this list.",
                "screenshot_path": asset_path("doctor-portal-update-card-and-profile", "02-doctor-select-child.png"),
                "caption": "Doctor-facing child selection after WhatsApp lookup",
                "what_to_show": "The list of matched children and their action buttons.",
                "require_asset": True,
            },
            {
                "title": "Open and update the vaccination card",
                "user_does": "Open the due-only or full vaccination card and record given dates where needed.",
                "user_sees": "The doctor-facing vaccination card with due, overdue, given, and future status states.",
                "why": "This is the operational heart of the clinic workflow: maintaining the child’s vaccination record.",
                "result": "The card reflects the latest known vaccination dates.",
                "notes": "The product automatically recalculates dependent doses when a given date changes.",
                "screenshot_path": asset_path("doctor-portal-update-card-and-profile", "03-doctor-vaccination-card.png"),
                "caption": "Doctor view of the patient vaccination card",
                "what_to_show": "The card layout with editable date inputs and schedule statuses.",
                "require_asset": True,
            },
            {
                "title": "Maintain the clinic profile",
                "user_does": "Open Edit doctor/clinic profile from the portal home.",
                "user_sees": "A profile form for doctor details, clinic contact data, languages, and IMC number.",
                "why": "Clinic identity, language settings, and contact details need to stay current for patient communication to work well.",
                "result": "The clinic profile can be updated without leaving the portal.",
                "notes": "Preferred languages influence downstream patient-facing messaging.",
                "screenshot_path": asset_path("doctor-portal-update-card-and-profile", "04-doctor-profile.png"),
                "caption": "Doctor and clinic profile maintenance screen",
                "what_to_show": "The profile form with clinic and doctor data.",
                "require_asset": True,
            },
        ],
    },
    {
        "number": 7,
        "group": "Doctor / Clinic",
        "slug": "reminders-and-education-workflows",
        "title": "Reminders and Education Workflows",
        "purpose": "Train clinic users on monitoring due doses, sending reminders, and opening the supporting vaccine education pages for providers and parents.",
        "primary_user": "Doctor or clinic staff member managing patient follow-up",
        "entry_point": "Doctor portal -> Reminders, child reminder screens, and education pages",
        "summary": "The clinic can review due and overdue doses, open child-level reminder schedules, and use linked education pages to reinforce communication with families. Reminder sending currently hands off to WhatsApp with a prefilled message.",
        "success_criteria": [
            "The clinic can read the reminder dashboard filters and eligible actions.",
            "Staff can inspect reminder status for one child.",
            "Provider and patient education pages can be opened from the workflow.",
        ],
        "related": [
            "06-doctor-portal-update-card-and-profile.md",
            "08-parent-self-service-card-history-and-share-link.md",
        ],
        "notes": [
            "Reminder sending opens WhatsApp rather than pushing through a backend messaging integration.",
            "The dashboard combines due-today / overdue logic with a retention window on recently overdue items.",
        ],
        "steps": [
            {
                "title": "Review the clinic-wide reminder dashboard",
                "user_does": "Open Reminders from the portal home.",
                "user_sees": "A reminder dashboard with status filters, vaccine filter, search, and Send Reminder actions.",
                "why": "This is the clinic’s queue view for follow-up work across all patients.",
                "result": "The clinic can identify who is due, overdue, upcoming, or already reminded.",
                "notes": "Search accepts child text and mobile-number lookups.",
                "screenshot_path": asset_path("reminders-and-education-workflows", "01-reminders-dashboard.png"),
                "caption": "Clinic-wide reminders dashboard",
                "what_to_show": "The dashboard filters, patient rows, status text, and reminder buttons.",
                "require_asset": True,
            },
            {
                "title": "Inspect one child’s reminder schedule",
                "user_does": "Open the reminder detail screen for a specific child.",
                "user_sees": "A child-level table showing vaccine reminder status and eligibility.",
                "why": "This narrows the follow-up conversation to one patient when the clinic needs detail.",
                "result": "The clinic can see exactly which reminder actions are available for the child.",
                "notes": "The child-level page is especially useful before sending a reminder manually.",
                "screenshot_path": asset_path("reminders-and-education-workflows", "02-child-reminders.png"),
                "caption": "Child-level reminder schedule in the doctor portal",
                "what_to_show": "The list of vaccine statuses and reminder actions for one child.",
                "require_asset": True,
            },
            {
                "title": "Open doctor-facing vaccine education",
                "user_does": "Open a vaccine education page from the reminder or card workflow.",
                "user_sees": "The provider education page for that vaccine.",
                "why": "Doctors may need refresher content or supporting material before discussing the vaccine with families.",
                "result": "The provider can access vaccine education without leaving the workflow context.",
                "notes": "This page is provider-focused, separate from the simplified patient education page used in reminder links.",
                "screenshot_path": asset_path("reminders-and-education-workflows", "03-doctor-vaccine-education.png"),
                "caption": "Doctor-facing vaccine education page",
                "what_to_show": "The provider education view opened from the workflow.",
                "require_asset": True,
            },
            {
                "title": "Understand the patient education destination",
                "user_does": "Open the simplified patient education page used in reminder messaging.",
                "user_sees": "A patient-friendly video page with language-aware vaccine education content.",
                "why": "This is the educational destination parents reach from reminder messages.",
                "result": "The clinic understands what families see after clicking an education link.",
                "notes": "This page is public and optimized for parent-facing education consumption.",
                "screenshot_path": asset_path("reminders-and-education-workflows", "04-patient-education-simple.png"),
                "caption": "Simplified patient education page used by reminder links",
                "what_to_show": "The public patient education view with vaccine videos.",
                "require_asset": True,
            },
        ],
    },
    {
        "number": 8,
        "group": "Parent / Patient",
        "slug": "parent-self-service-card-history-and-share-link",
        "title": "Parent Self-Service Card, History, and Share Link",
        "purpose": "Show how parents add or retrieve a child record, open the vaccination card, review history, use shared links, and access education content.",
        "primary_user": "Parent or guardian",
        "entry_point": "Direct parent routes `/add/`, `/update/`, and shared link `/p/<token>/`",
        "summary": "Parents can add a child, look up existing records by WhatsApp number, open the vaccination card, review vaccination history, verify a shared link, and open the patient education view. These flows are active in the product even though they are not promoted on the current public home page.",
        "success_criteria": [
            "The parent can reach the correct entry route for self-service.",
            "The parent can identify a child record and open the vaccination card.",
            "The parent can verify a share link and review history or education content.",
        ],
        "related": [
            "05-doctor-portal-add-patient-and-share-link.md",
            "07-reminders-and-education-workflows.md",
        ],
        "notes": [
            "Parent workflows are live but currently accessed through direct URLs rather than the public home page.",
        ],
        "steps": [
            {
                "title": "Use the public add-record route when starting fresh",
                "user_does": "Open the public Add Record page and review the child + parent fields.",
                "user_sees": "A simple add-record form for child details and the parent WhatsApp number.",
                "why": "This is the parent-side entry path when the family is not yet in the system.",
                "result": "The parent can create a new child record without using the doctor portal.",
                "notes": "This route now works in the local demo build after aligning the view with the live model fields.",
                "screenshot_path": asset_path("parent-self-service-card-history-and-share-link", "01-parent-add-form.png"),
                "caption": "Parent add-record form",
                "what_to_show": "The public add form for child details and parent WhatsApp number.",
                "require_asset": True,
            },
            {
                "title": "Look up existing children by WhatsApp number",
                "user_does": "Open Update Record, enter the parent WhatsApp number, and proceed.",
                "user_sees": "A lookup form followed by the list of matching children for that number.",
                "why": "This is the fastest way for a returning family to re-enter the system.",
                "result": "The parent can choose the correct child record.",
                "notes": "One parent number can return multiple children.",
                "screenshot_path": asset_path("parent-self-service-card-history-and-share-link", "03-parent-select-child.png"),
                "caption": "Parent child-selection screen after WhatsApp lookup",
                "what_to_show": "The child list returned after the parent WhatsApp lookup.",
                "require_asset": True,
            },
            {
                "title": "Open the vaccination card and history",
                "user_does": "Open the child vaccination card and then review the vaccination history page.",
                "user_sees": "The due-only vaccination card plus a separate history page with status badges and clinic call actions.",
                "why": "These are the parent’s main operational views for monitoring and reviewing vaccination progress.",
                "result": "The parent can see upcoming/due items and the historical record.",
                "notes": "The history page also links into education content.",
                "screenshot_path": asset_path("parent-self-service-card-history-and-share-link", "05-parent-history.png"),
                "caption": "Parent vaccination history view",
                "what_to_show": "The parent history page after a valid child session has been established.",
                "require_asset": True,
            },
            {
                "title": "Use a doctor-issued share link",
                "user_does": "Open the shared link from WhatsApp, verify the WhatsApp number, and continue into the child card.",
                "user_sees": "A verification page followed by the child vaccination card once the number matches.",
                "why": "This is how clinics hand off the card to parents without asking them to search manually.",
                "result": "The parent reaches the correct child record through the tokenized link.",
                "notes": "The shared link stores the expected last ten digits and asks the parent to confirm the number.",
                "screenshot_path": asset_path("parent-self-service-card-history-and-share-link", "06-share-link-verify.png"),
                "caption": "Parent share-link verification screen",
                "what_to_show": "The verification form asking the parent to confirm the WhatsApp number.",
                "require_asset": True,
            },
            {
                "title": "Open parent-facing vaccine education",
                "user_does": "Open the parent education link from the history or card workflow.",
                "user_sees": "A parent-facing education page with available vaccine video content.",
                "why": "Education supports understanding and follow-through, especially after reminders or history review.",
                "result": "The parent can move from record review into learning content without leaving the product journey.",
                "notes": "This uses the patient education library and language preference logic where content exists.",
                "screenshot_path": asset_path("parent-self-service-card-history-and-share-link", "08-parent-vaccine-education.png"),
                "caption": "Parent-facing vaccine education page",
                "what_to_show": "The vaccine education screen reached from the parent flow.",
                "require_asset": True,
            },
        ],
    },
]


def markdown_filename(workflow: dict[str, Any]) -> str:
    return f"{workflow['number']:02d}-{workflow['slug']}.md"


def deck_filename(workflow: dict[str, Any]) -> str:
    return f"{workflow['number']:02d}-{workflow['slug']}.pptx"


def add_textbox(slide, left, top, width, height, text="", *, font_size=20, bold=False, color=None, fill=None, line=None, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return shape


def add_title(slide, title: str, subtitle: str = "") -> None:
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(9.4), Inches(0.55), title, font_size=28, bold=True, color=COLORS["ink"])
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.06)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = COLORS["teal"]
    slide.shapes[-1].line.fill.background()
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.08), Inches(11.8), Inches(0.4), subtitle, font_size=12, color=COLORS["muted"])


def add_footer(slide, left_text: str, right_text: str) -> None:
    add_textbox(slide, Inches(0.6), Inches(7.05), Inches(6.0), Inches(0.3), left_text, font_size=10, color=COLORS["muted"])
    add_textbox(slide, Inches(9.8), Inches(7.05), Inches(2.9), Inches(0.3), right_text, font_size=10, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_image(slide, image_path: Path, left, top, width, height) -> None:
    if not image_path.exists():
        add_textbox(slide, left, top, width, height, f"Image missing:\n{image_path.name}", font_size=16, color=COLORS["muted"], fill=COLORS["surface"], line=COLORS["line"])
        return

    with Image.open(image_path) as img:
        img_w, img_h = img.size
    target_ratio = width / height
    image_ratio = img_w / img_h

    if image_ratio > target_ratio:
        final_h = height
        final_w = height * image_ratio
    else:
        final_w = width
        final_h = width / image_ratio

    pic = slide.shapes.add_picture(str(image_path), left, top, width=final_w, height=final_h)
    if final_w > width:
        pic.left = int(left - (final_w - width) / 2)
    if final_h > height:
        pic.top = int(top - (final_h - height) / 2)


def cover_slide(prs: Presentation, workflow: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["surface"]
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.25)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = COLORS["teal"]
    slide.shapes[-1].line.fill.background()
    add_textbox(slide, Inches(0.75), Inches(0.38), Inches(7.0), Inches(0.4), f"{workflow['number']:02d}", font_size=16, bold=True, color=RGBColor(255, 255, 255))
    add_textbox(slide, Inches(0.75), Inches(1.6), Inches(7.2), Inches(1.4), workflow["title"], font_size=28, bold=True, color=COLORS["ink"])
    add_textbox(slide, Inches(0.75), Inches(3.05), Inches(5.8), Inches(0.5), workflow["group"], font_size=13, bold=True, color=COLORS["teal"], fill=COLORS["teal_light"])
    add_textbox(slide, Inches(0.75), Inches(3.85), Inches(5.8), Inches(1.0), workflow["summary"], font_size=15, color=COLORS["muted"])
    add_textbox(slide, Inches(0.75), Inches(5.1), Inches(5.8), Inches(0.35), f"Primary user: {workflow['primary_user']}", font_size=12, color=COLORS["ink"])
    add_textbox(slide, Inches(0.75), Inches(5.45), Inches(5.8), Inches(0.35), f"Entry point: {workflow['entry_point']}", font_size=12, color=COLORS["ink"])
    image_step = next((s for s in workflow["steps"] if s.get("require_asset")), None)
    if image_step:
        add_image(slide, ROOT / image_step["screenshot_path"], Inches(7.0), Inches(1.4), Inches(5.6), Inches(4.7))
    add_footer(slide, "User-Flow Training Pack", workflow["slug"])


def summary_slide(prs: Presentation, workflow: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Workflow At A Glance", workflow["title"])
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(4.0), Inches(0.35), "What This Workflow Covers", font_size=16, bold=True, color=COLORS["teal"])
    add_textbox(slide, Inches(0.7), Inches(1.9), Inches(5.0), Inches(1.15), workflow["summary"], font_size=14, color=COLORS["ink"], fill=COLORS["surface"], line=COLORS["line"])
    add_textbox(slide, Inches(0.7), Inches(3.2), Inches(2.1), Inches(0.32), "Primary User", font_size=12, bold=True, color=COLORS["muted"])
    add_textbox(slide, Inches(0.7), Inches(3.52), Inches(5.0), Inches(0.42), workflow["primary_user"], font_size=14, color=COLORS["ink"])
    add_textbox(slide, Inches(0.7), Inches(4.05), Inches(2.1), Inches(0.32), "Entry Point", font_size=12, bold=True, color=COLORS["muted"])
    add_textbox(slide, Inches(0.7), Inches(4.37), Inches(5.0), Inches(0.55), workflow["entry_point"], font_size=14, color=COLORS["ink"])
    add_textbox(slide, Inches(0.7), Inches(5.05), Inches(3.0), Inches(0.32), "Success Criteria", font_size=12, bold=True, color=COLORS["muted"])
    box = slide.shapes.add_textbox(Inches(0.7), Inches(5.37), Inches(5.2), Inches(1.2))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(workflow["success_criteria"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(13)
        p.font.name = "Aptos"
        p.font.color.rgb = COLORS["ink"]
        p.bullet = True
    image_step = next((s for s in workflow["steps"] if s.get("require_asset")), None)
    if image_step:
        add_image(slide, ROOT / image_step["screenshot_path"], Inches(6.25), Inches(1.55), Inches(6.25), Inches(4.65))
    if workflow["number"] == 1:
        add_textbox(slide, Inches(6.25), Inches(6.32), Inches(6.25), Inches(0.28), "Role Map", font_size=12, bold=True, color=COLORS["muted"])
        add_textbox(slide, Inches(6.25), Inches(6.58), Inches(6.25), Inches(0.55), "Partner/Admin -> Doctor/Clinic -> Parent/Patient -> Reminders & Education", font_size=14, bold=True, color=COLORS["teal"], fill=COLORS["teal_light"])
    add_footer(slide, "Workflow summary", workflow["title"])


def step_slide(prs: Presentation, workflow: dict[str, Any], index: int, step: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, f"Step {index}. {step['title']}", workflow["title"])
    add_textbox(slide, Inches(0.7), Inches(1.55), Inches(4.8), Inches(0.28), "What the user does", font_size=12, bold=True, color=COLORS["teal"])
    add_textbox(slide, Inches(0.7), Inches(1.82), Inches(5.0), Inches(0.85), step["user_does"], font_size=14, color=COLORS["ink"], fill=COLORS["surface"], line=COLORS["line"])
    add_textbox(slide, Inches(0.7), Inches(2.83), Inches(4.8), Inches(0.28), "What the user sees", font_size=12, bold=True, color=COLORS["teal"])
    add_textbox(slide, Inches(0.7), Inches(3.1), Inches(5.0), Inches(0.88), step["user_sees"], font_size=14, color=COLORS["ink"], fill=COLORS["surface"], line=COLORS["line"])
    add_textbox(slide, Inches(0.7), Inches(4.16), Inches(4.8), Inches(0.28), "Why this matters", font_size=12, bold=True, color=COLORS["coral"])
    add_textbox(slide, Inches(0.7), Inches(4.43), Inches(5.0), Inches(0.82), step["why"], font_size=14, color=COLORS["ink"], fill=COLORS["coral_light"], line=COLORS["line"])
    add_textbox(slide, Inches(0.7), Inches(5.43), Inches(4.8), Inches(0.28), "Expected result", font_size=12, bold=True, color=COLORS["teal"])
    add_textbox(slide, Inches(0.7), Inches(5.7), Inches(5.0), Inches(0.7), step["result"], font_size=14, color=COLORS["ink"], fill=COLORS["surface"], line=COLORS["line"])
    if step.get("notes"):
        add_textbox(slide, Inches(0.7), Inches(6.48), Inches(5.0), Inches(0.5), f"Trainer note: {step['notes']}", font_size=11, color=COLORS["muted"])
    add_image(slide, ROOT / step["screenshot_path"], Inches(6.0), Inches(1.55), Inches(6.55), Inches(5.15))
    add_textbox(slide, Inches(6.05), Inches(6.75), Inches(6.45), Inches(0.28), step["caption"], font_size=11, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_footer(slide, f"Step {index} of {len(workflow['steps'])}", workflow["title"])


def notes_slide(prs: Presentation, workflow: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Common Issues and Trainer Tips", workflow["title"])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.6), Inches(4.8))
    tf = box.text_frame
    tf.clear()
    items = workflow.get("notes", [])
    if not items:
        items = ["No additional trainer notes were captured for this workflow."]
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = "Aptos"
        p.font.color.rgb = COLORS["ink"]
        p.bullet = True
    add_textbox(slide, Inches(0.8), Inches(6.15), Inches(11.6), Inches(0.55), "Use these notes to pre-empt confusion before the audience encounters the step live.", font_size=12, color=COLORS["muted"], fill=COLORS["surface"], line=COLORS["line"])
    add_footer(slide, "Trainer guidance", workflow["title"])


def closing_slide(prs: Presentation, workflow: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["teal"]
    add_textbox(slide, Inches(0.9), Inches(1.2), Inches(8.5), Inches(0.65), workflow["title"], font_size=28, bold=True, color=RGBColor(255, 255, 255))
    add_textbox(slide, Inches(0.9), Inches(2.0), Inches(8.8), Inches(0.45), "Related Documents", font_size=14, bold=True, color=RGBColor(255, 255, 255))
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.45), Inches(7.4), Inches(2.5))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(workflow["related"]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.name = "Aptos"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.bullet = True
    add_textbox(slide, Inches(0.9), Inches(5.55), Inches(10.8), Inches(0.55), workflow["purpose"], font_size=14, color=RGBColor(255, 255, 255))
    add_textbox(slide, Inches(0.9), Inches(6.25), Inches(10.8), Inches(0.45), STATUS_TEXT, font_size=11, color=RGBColor(255, 255, 255))
    add_footer(slide, "Training pack close", workflow["slug"])


def create_workflow_deck(workflow: dict[str, Any]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover_slide(prs, workflow)
    summary_slide(prs, workflow)
    for idx, step in enumerate(workflow["steps"], start=1):
        step_slide(prs, workflow, idx, step)
    notes_slide(prs, workflow)
    closing_slide(prs, workflow)
    deck_path = DECK_DIR / deck_filename(workflow)
    prs.save(deck_path)
    return deck_path


def create_master_index(workflows: list[dict[str, Any]]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = COLORS["surface"]
    cover.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15)).fill.solid()
    cover.shapes[-1].fill.fore_color.rgb = COLORS["teal"]
    cover.shapes[-1].line.fill.background()
    add_textbox(cover, Inches(0.75), Inches(1.6), Inches(7.8), Inches(0.9), "User-Flow Training Index", font_size=30, bold=True, color=COLORS["ink"])
    add_textbox(cover, Inches(0.75), Inches(2.6), Inches(9.6), Inches(0.8), "Open this deck first. Each link below opens a sibling workflow deck in the shared package.", font_size=16, color=COLORS["muted"])
    add_textbox(cover, Inches(0.75), Inches(4.2), Inches(5.2), Inches(0.42), "Sections included", font_size=14, bold=True, color=COLORS["teal"])
    add_textbox(cover, Inches(0.75), Inches(4.7), Inches(5.5), Inches(1.4), "Platform overview\nAdmin / Partner\nDoctor / Clinic\nParent / Patient", font_size=18, color=COLORS["ink"], fill=COLORS["surface"], line=COLORS["line"])
    add_footer(cover, "Master index deck", "00-user-flow-training-index.pptx")

    grouped: dict[str, list[dict[str, Any]]] = {}
    for workflow in workflows:
        grouped.setdefault(workflow["group"], []).append(workflow)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Workflow Deck Directory", "Each card links to a sibling .pptx file in this folder.")
    y = 1.45
    for group, items in grouped.items():
        add_textbox(slide, Inches(0.75), Inches(y), Inches(3.5), Inches(0.32), group, font_size=14, bold=True, color=COLORS["teal"])
        y += 0.45
        for item in items:
            box = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(11.3), Inches(0.42))
            box.fill.solid()
            box.fill.fore_color.rgb = COLORS["surface"]
            box.line.color.rgb = COLORS["line"]
            p = box.text_frame.paragraphs[0]
            p.text = f"{item['number']:02d}. {item['title']}"
            run = p.runs[0]
            run.font.name = "Aptos"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = COLORS["ink"]
            run.hyperlink.address = deck_filename(item)
            y += 0.55
        y += 0.18
    add_footer(slide, "Master index", "Click a title to open the deck")

    index_path = DECK_DIR / "00-user-flow-training-index.pptx"
    prs.save(index_path)
    return index_path


def write_manual(workflow: dict[str, Any]) -> Path:
    path = MANUAL_DIR / markdown_filename(workflow)
    lines = [
        f"# {workflow['title']}",
        "",
        "## 1. Title",
        workflow["title"],
        "",
        "## 2. Document Purpose",
        workflow["purpose"],
        "",
        "## 3. Primary User",
        workflow["primary_user"],
        "",
        "## 4. Entry Point",
        workflow["entry_point"],
        "",
        "## 5. Workflow Summary",
        workflow["summary"],
    ]
    if workflow.get("diagram"):
        lines.extend(["", workflow["diagram"]])
    lines.extend(["", "## 6. Step-By-Step Instructions", ""])
    for idx, step in enumerate(workflow["steps"], start=1):
        lines.extend(
            [
                f"### Step {idx}. {step['title']}",
                f"- What the user does: {step['user_does']}",
                f"- What the user sees: {step['user_sees']}",
                f"- Why the step matters: {step['why']}",
                f"- Expected result: {step['result']}",
                f"- Common issues or trainer notes: {step.get('notes', 'None noted.')}",
                "- Screenshot placeholder section:",
                f"  - Suggested file path: `{step['screenshot_path']}`",
                f"  - Screenshot caption: {step['caption']}",
                f"  - What the screenshot should show: {step['what_to_show']}",
                "",
            ]
        )
    lines.extend(["## 7. Success Criteria"])
    for item in workflow["success_criteria"]:
        lines.append(f"- {item}")
    lines.extend(["", "## 8. Related Documents"])
    for rel in workflow["related"]:
        lines.append(f"- [{rel}]({rel})")
    lines.extend(["", "## 9. Status", STATUS_TEXT, ""])
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def verify_links(master_index: Path, workflows: list[dict[str, Any]]) -> list[str]:
    prs = Presentation(str(master_index))
    expected = {deck_filename(workflow) for workflow in workflows}
    seen: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    address = run.hyperlink.address
                    if address:
                        seen.add(address)
    missing = sorted(expected - seen)
    return missing


def verify_package(workflows: list[dict[str, Any]], master_index: Path) -> str:
    issues: list[str] = []
    for workflow in workflows:
        for step in workflow["steps"]:
            if step.get("require_asset") and not (ROOT / step["screenshot_path"]).exists():
                issues.append(f"Missing screenshot asset: {step['screenshot_path']}")
    placeholder_pattern = re.compile(r"\bTODO\b|\bTBD\b|lorem ipsum", re.IGNORECASE)
    for deck in DECK_DIR.glob("*.pptx"):
        with zipfile.ZipFile(deck) as zf:
            for name in zf.namelist():
                if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
                    continue
                text = zf.read(name).decode("utf-8", errors="ignore")
                if placeholder_pattern.search(text):
                    issues.append(f"Placeholder text found in {deck.name}::{name}")
    missing_links = verify_links(master_index, workflows)
    for missing in missing_links:
        issues.append(f"Master index missing hyperlink for {missing}")

    report_lines = [
        "User-flow training pack QA report",
        "================================",
        f"Workflow manuals: {len(list(MANUAL_DIR.glob('*.md')))}",
        f"Workflow decks: {len(list(DECK_DIR.glob('[0-9][0-9]-*.pptx')))}",
        f"Master index: {master_index.name}",
        "",
    ]
    if issues:
        report_lines.append("Issues found:")
        report_lines.extend(f"- {issue}" for issue in issues)
    else:
        report_lines.append("No missing screenshots, placeholder strings, or master-index link gaps detected.")
    report = "\n".join(report_lines) + "\n"
    QA_REPORT.write_text(report, encoding="utf-8")
    return report


def main() -> None:
    MANUAL_DIR.mkdir(parents=True, exist_ok=True)
    DECK_DIR.mkdir(parents=True, exist_ok=True)

    for workflow in WORKFLOWS:
        write_manual(workflow)
        create_workflow_deck(workflow)

    master_index = create_master_index(WORKFLOWS)
    report = verify_package(WORKFLOWS, master_index)
    print(report)


if __name__ == "__main__":
    main()
